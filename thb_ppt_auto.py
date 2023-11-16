# -*- coding: utf-8 -*-
"""
Created on Mon Nov  6 17:39:45 2023

@author: User
"""

import os
import requests
from selenium import webdriver
from selenium.webdriver.common.by import By
from io import BytesIO
from PIL import Image
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re

# 測試 1234565756 
#測試123456123456
hsdaoasidhaslkdhkasd

# 圖片儲存位置
img_save_path = r"C:\Users\User\Desktop\ppt_自動化\result"
os.chdir(img_save_path)

now = datetime.now()
one_day = timedelta(days = 1)
one_hour = timedelta(hours = 1)
yesterday = now - one_day
next_day = now + one_day
day_after_tomorrow = now + 2*one_day
one_hours_ago = now - one_hour

# webdriver
driver = webdriver.Chrome()

# 擷取氣象局圖資網址 (除了QPF)
def get_image_url(url, location, item):
    url_list = []
    try:
        driver.get(url)
        for i in location :
            img_url = driver.find_elements(By.CLASS_NAME, "img-responsive")[i].get_attribute("src")
            url_list.append(img_url)
    except Exception as e:
        print(f"獲取{item}網址時發生未預期的錯誤：{e}")
    finally:
        return url_list

# 儲存成png檔案
def fetch_image(url, image_name):
    response = requests.get(url)
    if response.status_code == 200 :
        with Image.open(BytesIO(response.content)) as image:
            image.save(os.path.join(img_save_path, f"{image_name}.png"))
    elif url == fr"http://tidpweather.vpnplus.to:50020/RainMap/Taiwan/TW_{now.year}{now.month}{now.day}_{now.hour}00.png" and response.status_code != 200:
           print(f"RainMap未更新{now.hour}累積雨量，採用前一小時資料")
           E_url = rf"http://tidpweather.vpnplus.to:50020/RainMap/Taiwan/TW_{now.year}{now.month}{now.day}_{one_hours_ago.hour}00.png"
           fetch_image(E_url, "E_image") 
    else :
         print(f"{image_name} 圖資未更新或是發生預期外錯誤, 錯誤代碼: {response.status_code}")

# 擷取QPF圖資網址
def QPF_fetch_image(QPF_url):
    QPF_list = []
    for url in QPF_url:
        pattern = r"_ChFcstPrecip_([\w_]+)\.png"
        match = re.search(pattern, url)
        img_name = match[1]
        try:
            response = requests.get(url)
            QPF_image = Image.open(BytesIO(response.content))
            QPF_image.save(f"{img_name}QPF.png")
            QPF_list.append(QPF_image)
        except Exception as e:
            print(f"QPF圖資載入時發生錯誤，錯誤代碼如下：{e}")
    return QPF_list

# 結合QPF圖資
def combine_image(QPF_list, forecast_range):
    bg = Image.new("RGBA", (2490, 3000), "#FFFFFF")
    for i in range(1,5):
        image = QPF_list[i-1]
        x = (i-1) % 2
        y = (i-1) // 2
        bg.paste(image, (x*1245, y*1500))
    bg.save(f"{forecast_range}hr_total.png")

# 地面天氣圖
SWM_url = get_image_url("https://www.cwa.gov.tw/V8/C/W/analysis.html", [2], "地面天氣圖")
fetch_image(SWM_url[0], "SWM")

# 雷達迴波圖
Radar_url = get_image_url("https://www.cwa.gov.tw/V8/C/W/OBS_Radar.html", [1], "雷達迴波圖")
fetch_image(Radar_url[0], "Radar")

# 流場圖
StreamLine_url = "https://tropic.ssec.wisc.edu/real-time/westpac/winds/wgmsdlm1.GIF"
fetch_image(StreamLine_url, "StreamLine")

# 衛星雲圖
driver.get(r"https://www.cwa.gov.tw/V8/C/W/OBS_Sat.html?Area=2")
button = driver.find_element(By.LINK_TEXT, "真實色")
button.click()
Satellite_Images_URL = driver.find_elements(By.CLASS_NAME, "img-responsive")[1].get_attribute("src")
fetch_image(Satellite_Images_URL, "Satellite")

# CWA定量降水網址
location_range = list(range(2,10)) # list中第2~10是QPF網址
QPF_url = get_image_url("https://www.cwa.gov.tw/V8/C/P/QPF.html", location_range, "QPF")
QPF_list = QPF_fetch_image(QPF_url)
QPF_6hr = QPF_list[4:]
QPF_12hr = QPF_list[:4]
combine_image(QPF_6hr, 6)
combine_image(QPF_12hr, 12)

driver.close()


# 公司_今日累積雨量圖資
E_url = rf"http://tidpweather.vpnplus.to:50020/RainMap/Taiwan/TW_{now.year}{now.month}{now.day}_{now.hour}00.png"
fetch_image(E_url, "E_image")

# 公司_每日06時累積雨量
day = now.strftime("%Y%m%d")
E_url = fr"http://tidpweather.vpnplus.to:50020/RainMap/Taiwan/TW_{now.year}{now.month}{now.day}_0600.png"
fetch_image(E_url, "E_06_image")

# 公司_昨日累積雨量
E_url = fr"http://tidpweather.vpnplus.to:50020/RainMap/Taiwan/TW_{now.year}{now.month}{now.day}_0000.png"
fetch_image(E_url, "E_yday_image")


#==============================================================================


try:
    prs = Presentation("本局-20231113 未來三日天氣分析報告.pptx")
except Exception as e:
    print(e)
    
def convert_to_minguo(year):
    return year - 1911

# 字型設定
def font(textbox, Date_String, RGB, size):
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.clear()
    run = paragraph.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = RGB
    run.text = Date_String
    run.font.bold = True
    run.font.name = "微軟正黑體"
    paragraph.alignment = PP_ALIGN.CENTER

# 更新文字方塊日期
def update_date(slide, re_express, Date_String, RGB, size):
    pattern = re.compile(re_express)
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                raw_text = shape.text_frame.text
                # is "raw_text" match pattern
                if re.search(pattern, raw_text):
                    font(shape, Date_String, RGB, size)
                    break 
        except Exception as e:
            print(f"{slide} update date have unexpected exception: {e}")

# 更新表格內日期
def table_update_date(slide, start, end, RGB, size):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for i in range(start, end):
                slide_text_box = shape.table.cell(0, i)
                if i == start:
                    Date_String = f"今({now.day})日"
                elif i == (start + 1):
                    Date_String = f"明({next_day.day})日"
                elif i == (start + 2):
                    Date_String = f"後({day_after_tomorrow.day})日"
                font(slide_text_box, Date_String, RGB, size)


first_page_pattern = r"時間：\d{3}年\d{1,2}月\d{1,2}日 \d{4}時"
SWM_pattern = r"\d{1,2}月\d{1,2}日 02:00 地面天氣圖"
Satellite_pattern = r"\d{1,2}月\d{1,2}日 \d{1,2}:00 衛星雲圖"
StreamLine_pattern = r"\d{1,2}月\d{1,2}日 05:00 700-850hPa 平均駛流場圖"
yday_accumulate_pattern = r"昨\(\d{1,2}\)日 累積雨量"
tday_accumulate_pattern = r"今\(\d{1,2}\)日 00-\d{2}時 累積雨量"
today_QPF_pattern = r"今\(\d{1,2}\)日$"
tomorrow_QPF_pattern = r"明\(\d{1,2}\)日$"
day_after_tomorrow_QPF_pattern = r"後\(\d{1,2}\)日$"

first_slide = prs.slides[0]
second_slide = prs.slides[1]
thrid_slide = prs.slides[2]
forth_slide = prs.slides[3]
fifth_slide = prs.slides[4]
sixth_slide = prs.slides[5]
seventh_slide = prs.slides[6]

update_date(first_slide, first_page_pattern, f"{convert_to_minguo(now.year)}年{now.month}月{now.day}日 09:00時", RGBColor(00,00,00), 24)
update_date(second_slide, SWM_pattern, f"{now.month}月{now.day}日 02:00 地面天氣圖", RGBColor(12,51,115), 18)
update_date(second_slide, Satellite_pattern, f"{now.month}月{now.day}日 07:00 衛星雲圖", RGBColor(12,51,115), 18)
update_date(thrid_slide, StreamLine_pattern, f"{now.month}月{now.day}日 05:00 700-850hPa 平均駛流場圖", RGBColor(12,51,115), 18)
update_date(forth_slide, yday_accumulate_pattern, f"昨({yesterday.day})日 累積雨量", RGBColor(12,51,115), 18)
update_date(forth_slide, tday_accumulate_pattern, f"今({now.day})日 00-06時 累積雨量", RGBColor(12,51,115), 18)
update_date(fifth_slide, today_QPF_pattern, f"今({now.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, tomorrow_QPF_pattern, f"明({next_day.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, day_after_tomorrow_QPF_pattern, f"今({day_after_tomorrow.day})日", RGBColor(12,51,115), 18)
table_update_date(fifth_slide, 1, 4, RGBColor(00,00,00), 18)
table_update_date(sixth_slide, 3, 6, RGBColor(00,00,00), 18)
table_update_date(seventh_slide, 3, 6, RGBColor(00,00,00), 18)


# 更換圖資
def change_img(slide, img, left):
    for shape in slide.shapes:
        if shape.shape_type == 13 and (left - 50000<= shape.left <= left + 50000): # 13代表圖片
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # 刪除圖片
            sp = shape._element
            sp.getparent().remove(sp)
            # 插入新圖片
            slide.shapes.add_picture(img,  left, top, width, height)
            print(f"{img} success")
            print(img, left, top, width, height, "\n")
            break

change_img(second_slide, "Satellite.png", 6988336)
change_img(second_slide, "SWM.png", 1106797)   
change_img(thrid_slide, "StreamLine.png", 835086)
change_img(forth_slide, "E_06_image.png", 6816100)
change_img(forth_slide, "E_yday_image.png", 1403989)

format_date = now.strftime("%Y%m%d")
prs.save(f'本局-{format_date} 未來三日天氣分析報告.pptx')
