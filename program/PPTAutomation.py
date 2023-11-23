# -*- coding: utf-8 -*-
"""
Created on Fri Nov 17 17:15:19 2023

@author: User
"""


from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import re
import os
import gc


# 創建時間設定物件
now = datetime.now()
yesterday = now - timedelta(days = 1)
next_day = now + timedelta(days = 1)
day_after_tomorrow = now + timedelta(days = 2)


# 設定時間格式
today = now.strftime("%Y%m%d")
yday = yesterday.strftime("%Y%m%d")


# 設定路徑
path = rf"C:\Users\user\Desktop\自動化separate\THB_PPT_Automation\{today}"
yday_ppt_path = rf"C:\Users\user\Desktop\自動化separate\THB_PPT_Automation\{yday}\本局-{yday} 未來三日天氣分析報告.pptx"
os.chdir(path)
    
    

def write_txt(text):
    date = now.strftime("%y%m%d")
    with open(os.path.join(path, f"{date}_img.txt"), 'a') as file:
        file.write(text + "\n")
    
    
    
# 西元轉換民國
def convert_to_minguo(year):
    return year - 1911


# 字型設定
def font(textbox, Date_String, RGB, size):
    
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.clear()
    run = paragraph.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = RGB
    run.text = Date_String
    run.font.bold = True
    run.font.name = "微軟正黑體"
    paragraph.alignment = PP_ALIGN.CENTER

# 更新文字方塊日期
def update_date(slide, re_express, Date_String, RGB, size):
    pattern = re.compile(re_express)
    
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                raw_text = shape.text_frame.text
                # is "raw_text" match pattern
                if re.search(pattern, raw_text):
                    font(shape, Date_String, RGB, size)
                    break 
                
        except Exception as e:
            write_txt(f"{slide} update have unexpected exception: {e}")


# 更新表格內日期
def table_update_date(slide, start, end, RGB, size):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for i in range(start, end):
                slide_text_box = shape.table.cell(0, i)
                if i == start:
                    Date_String = f"今({now.day})日"
                elif i == (start + 1):
                    Date_String = f"明({next_day.day})日"
                elif i == (start + 2):
                    Date_String = f"後({day_after_tomorrow.day})日"
                font(slide_text_box, Date_String, RGB, size)
             
                
                
# 讀取簡報檔案
try:
    
    if os.path.exists(f'本局-{today} 未來三日天氣分析報告.pptx'):
        prs = Presentation(f'本局-{today} 未來三日天氣分析報告.pptx')
        write_txt("採用今日簡報")

    else:
        prs = Presentation(yday_ppt_path)
        write_txt("採用昨日簡報")

    
except Exception as e:
    
    write_txt(f"exception occur when open ppt, error：{e}")



# regular expression
first_page_pattern = r"\d{3}年\d{1,2}月\d{1,2}日 \d{4}時"
SWM_pattern = r"\d{1,2}月\d{1,2}日 02:00 地面天氣圖"
Satellite_pattern = r"\d{1,2}月\d{1,2}日 \d{1,2}:00 衛星雲圖"
StreamLine_pattern = r"\d{1,2}月\d{1,2}日 05:00 700-850hPa 平均駛流場圖"
yday_accumulate_pattern = r"昨\(\d{1,2}\)日 累積雨量"
tday_accumulate_pattern = r"今\(\d{1,2}\)日 00-\d{2}時 累積雨量"
today_QPF_pattern = r"今\(\d{1,2}\)日$"
tomorrow_QPF_pattern = r"明\(\d{1,2}\)日$"
day_after_tomorrow_QPF_pattern = r"後\(\d{1,2}\)日$"


first_slide = prs.slides[0]
second_slide = prs.slides[1]
thrid_slide = prs.slides[2]
forth_slide = prs.slides[3]
fifth_slide = prs.slides[4]
sixth_slide = prs.slides[5]
seventh_slide = prs.slides[6]

update_date(first_slide, first_page_pattern, f"{convert_to_minguo(now.year)}年{now.month}月{now.day}日 0900時", RGBColor(00,00,00), 24)
update_date(second_slide, SWM_pattern, f"{now.month}月{now.day}日 02:00 地面天氣圖", RGBColor(12,51,115), 18)
update_date(second_slide, Satellite_pattern, f"{now.month}月{now.day}日 07:00 衛星雲圖", RGBColor(12,51,115), 18)
update_date(thrid_slide, StreamLine_pattern, f"{now.month}月{now.day}日 05:00 700-850hPa 平均駛流場圖", RGBColor(12,51,115), 18)
update_date(forth_slide, yday_accumulate_pattern, f"昨({yesterday.day})日 累積雨量", RGBColor(12,51,115), 18)
update_date(forth_slide, tday_accumulate_pattern, f"今({now.day})日 00-06時 累積雨量", RGBColor(12,51,115), 18)
update_date(fifth_slide, today_QPF_pattern, f"今({now.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, tomorrow_QPF_pattern, f"明({next_day.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, day_after_tomorrow_QPF_pattern, f"後({day_after_tomorrow.day})日", RGBColor(12,51,115), 18)
table_update_date(fifth_slide, 1, 4, RGBColor(00,00,00), 18)
table_update_date(sixth_slide, 3, 6, RGBColor(00,00,00), 18)
table_update_date(seventh_slide, 3, 6, RGBColor(00,00,00), 18)




# 更換圖資
def change_img(slide, img, left):
    
    for shape in slide.shapes:
        if shape.shape_type == 13 and (left - 50000<= shape.left <= left + 50000): # 13代表圖片
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # 刪除圖片
            sp = shape._element
            sp.getparent().remove(sp)
            # 插入新圖片
            slide.shapes.add_picture(img,  left, top, width, height)
            
            write_txt(f"{img} success")
            write_txt((f"{img}, {left} \n"))
            break


change_img(second_slide, "round_Satellite.png", 6988336)
change_img(second_slide, "round_SWM.png", 1106797)   
change_img(thrid_slide, "round_StreamLine.png", 835086)
# change_img(forth_slide, "E_06_image.png", 6816100)
# change_img(forth_slide, "E_yday_image.png", 1403989)

try:
    
    prs.save(f'本局-{today} 未來三日天氣分析報告.pptx')
    write_txt("簡報存成功")

except Exception as e:
    write_txt(e)

finally:
    gc.collect()
