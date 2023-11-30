# -*- coding: utf-8 -*-
"""
Created on Fri Nov 17 17:15:19 2023

@author: User
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import pandas as pd 
import re
import os

now = datetime.now()
yesterday = now - timedelta(days = 1)
next_day = now + timedelta(days = 1)
day_after_tomorrow = now + timedelta(days = 2)
format_date = now.strftime("%Y%m%d")
format_ydate = yesterday.strftime("%Y%m%d")


# 程式執行位置
path = rf"C:\Users\User\Desktop\ppt_自動化\{format_date}"
# 今日檔案位置
ppt_path = rf"C:\Users\User\Desktop\ppt_自動化\{format_date}\本局-{format_date} 未來三日天氣分析報告.pptx"
# 昨日檔案位置
ppt_ypath = rf"C:\Users\User\Desktop\ppt_自動化\{format_ydate}\本局-{format_ydate} 未來三日天氣分析報告.pptx"
# 昨日雨量csv檔案位置
csv_path = rf"C:\Users\User\Desktop\ppt_自動化\yday_accumulate\{format_date}\{format_ydate}_累積雨量.csv"
os.chdir(path)

    
def convert_to_minguo(year):
    
    return year - 1911



def write_txt(text):
    date = now.strftime("%y%m%d")
    with open(os.path.join(path, f"{date}_img.txt"), 'a') as file:
        file.write(text + "\n")



# 字型設定
def font(textbox, Date_String, RGB, size):
    
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.clear()
    run = paragraph.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = RGB
    run.text = Date_String
    run.font.bold = True
    run.font.name = "微軟正黑體"
    paragraph.alignment = PP_ALIGN.CENTER

# 更新文字方塊日期
def update_date(slide, re_express, Date_String, RGB, size):
    pattern = re.compile(re_express)
    
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                raw_text = shape.text_frame.text
                # is "raw_text" match pattern
                if re.search(pattern, raw_text):
                    font(shape, Date_String, RGB, size)
                    break 
                
        except Exception as e:
            write_txt(f"{slide} update have unexpected exception: {e}")


# 更新表格內日期
def table_update_date(slide, start, end, RGB, size):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for i in range(start, end):
                slide_text_box = shape.table.cell(0, i)
                if i == start:
                    Date_String = f"今({now.day})日"
                elif i == (start + 1):
                    Date_String = f"明({next_day.day})日"
                elif i == (start + 2):
                    Date_String = f"後({day_after_tomorrow.day})日"
                font(slide_text_box, Date_String, RGB, size)
                
                
try:
    if os.path.exists(ppt_path):
        
        write_txt("採用今日簡報")
        prs = Presentation(ppt_path)
    
    else:
        
        write_txt("採用昨日簡報")
        prs = Presentation(ppt_ypath)
    
except Exception as e:
    write_txt(f"exception occur when open ppt, error：{e}")
    
    


first_page_pattern = r"\d{3}年\d{1,2}月\d{1,2}日 \d{4}時"
SWM_pattern = r"\d{1,2}月\d{1,2}日 02:00 地面天氣圖"
Satellite_pattern = r"\d{1,2}月\d{1,2}日 \d{1,2}:\d{1,2} 衛星雲圖"
StreamLine_pattern = r"\d{1,2}月\d{1,2}日 05:00 700-850hPa 平均駛流場圖"
yday_accumulate_pattern = r"昨\(\d{1,2}\)日 累積雨量"
tday_accumulate_pattern = r"今\(\d{1,2}\)日 00-\d{2}時 累積雨量"
today_QPF_pattern = r"今\(\d{1,2}\)日$"
tomorrow_QPF_pattern = r"明\(\d{1,2}\)日$"
day_after_tomorrow_QPF_pattern = r"後\(\d{1,2}\)日$"

first_slide = prs.slides[0]
second_slide = prs.slides[1]
thrid_slide = prs.slides[2]
forth_slide = prs.slides[3]
fifth_slide = prs.slides[4]
sixth_slide = prs.slides[5]
seventh_slide = prs.slides[6]

update_date(first_slide, first_page_pattern, f"{convert_to_minguo(now.year)}年{now.month}月{now.day}日 0900時", RGBColor(00,00,00), 24)
update_date(second_slide, SWM_pattern, f"{now.month}月{now.day}日 02:00 地面天氣圖", RGBColor(12,51,115), 18)
update_date(second_slide, Satellite_pattern, f"{now.month}月{now.day}日 07:00 衛星雲圖", RGBColor(12,51,115), 18)
update_date(thrid_slide, StreamLine_pattern, f"{now.month}月{now.day}日 05:00 700-850hPa 平均駛流場圖", RGBColor(12,51,115), 18)
update_date(forth_slide, yday_accumulate_pattern, f"昨({yesterday.day})日 累積雨量", RGBColor(12,51,115), 18)
update_date(forth_slide, tday_accumulate_pattern, f"今({now.day})日 00-06時 累積雨量", RGBColor(12,51,115), 18)
update_date(fifth_slide, today_QPF_pattern, f"今({now.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, tomorrow_QPF_pattern, f"明({next_day.day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, day_after_tomorrow_QPF_pattern, f"後({day_after_tomorrow.day})日", RGBColor(12,51,115), 18)
table_update_date(fifth_slide, 1, 4, RGBColor(00,00,00), 18)
table_update_date(sixth_slide, 3, 6, RGBColor(00,00,00), 18)
table_update_date(seventh_slide, 3, 6, RGBColor(00,00,00), 18)





# 更換圖資
def change_img(slide, img, left):
    
    for shape in slide.shapes:
        if shape.shape_type == 13 and (left - 50000<= shape.left <= left + 50000): # 13代表圖片
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # 刪除圖片
            sp = shape._element
            sp.getparent().remove(sp)

            slide.shapes.add_picture(img,  left, top, width, height)
            
            write_txt(f"{img} success  ")
            write_txt(f"{img}, {left}, {top}, {width}, {height} \n")
            break


change_img(second_slide, "round_Satellite.png", 6988336)
change_img(second_slide, "round_SWM.png", 1106797)   
change_img(thrid_slide, "round_StreamLine.png", 835086)
change_img(forth_slide, "E_06_image.png", 6816100)
change_img(forth_slide, "E_yday_image.png", 1403989)





#====================

# 回傳一個list
def read_csv(path):
    
    df = pd.read_csv(csv_path, encoding="big5")

    df.set_index("Unnamed: 0", inplace=True)

    df.index.name = "測站"
    
    df_dict = df.to_dict(orient="dict")[f"昨日({format_ydate})累積雨量"]
    
    rainfall_list = []
    
    for key, item in df_dict.items():
        
        rainfall_list.append(str(item))
    
    return rainfall_list



# 昨日累積雨量
def table_yday_rainfall(slide, start, end):
        
    for shape in slide.shapes:
        
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            
            count = 1
            
            for item in read_csv(csv_path)[start:end]:
                
                rainfall_box = shape.table.cell(count,2)
                font(rainfall_box, item, RGBColor(00,00,00), 16)
                count = count + 1
                

try:

    table_yday_rainfall(sixth_slide, 0, 10)
    table_yday_rainfall(seventh_slide,10, 21)

except:
    
    write_txt("未更新監控路段昨日雨量")


prs.save(f'本局-{format_date} 未來三日天氣分析報告.pptx')
